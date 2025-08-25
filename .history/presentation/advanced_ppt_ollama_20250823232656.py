from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import json
import os
from datetime import datetime

INPUT_FILE = "processing/cluster_summaries_ollama.json"
OUTPUT_DIR = "presentation"

def create_ai_ppt(summary_data):
    prs = Presentation()

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "📊 FT.com Daily Summary"
    slide1.placeholders[1].text = f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}"
    slide1.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)  # Blue

    # Slide 2: Themes & Articles
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "📝 Articles by Theme"
    tf2 = slide2.placeholders[1].text_frame
    tf2.clear()
    for item in summary_data["articles"]:
        theme = item.get("theme", "Miscellaneous")
        title = item.get("title", "")
        summary = item.get("summary", "")
        p = tf2.add_paragraph()
        p.text = f"{theme}: {title}"
        p.font.size = Pt(14)

    # Slide 3: Summaries (longer, auto split if needed)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "📝 Summaries"
    tf3 = slide3.placeholders[1].text_frame
    tf3.clear()
    for item in summary_data["articles"]:
        title = item.get("title", "")
        summary = item.get("summary", "")
        p = tf3.add_paragraph()
        p.text = f"{title}: {summary}"
        p.font.size = Pt(12)

    # Slide 4: References / Links
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "📰 Articles Included"
    tf4 = slide4.placeholders[1].text_frame
    tf4.clear()
    for item in summary_data["articles"]:
        title = item.get("title", "")
        link = item.get("link", "")
        p = tf4.add_paragraph()
        p.text = f"{title} ({link})"
        p.font.size = Pt(12)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_file = os.path.join(OUTPUT_DIR, f"FT_Summary_4Slides_{datetime.utcnow().date()}.pptx")
    prs.save(output_file)
    print(f"✅ Presentation saved: {output_file}")


if __name__ == "__main__":
    if not os.path.exists(INPUT_FILE):
        raise FileNotFoundError(f"❌ Input file not found: {INPUT_FILE}")

    with open(INPUT_FILE, "r", encoding="utf-8") as f:
        summary_data = json.load(f)

    create_ai_ppt(summary_data)
