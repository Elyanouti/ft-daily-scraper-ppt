import json
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

INPUT_FILE = "processing/cluster_summaries_ollama.json"
OUTPUT_DIR = "presentation"

MAX_PARAGRAPHS_PER_SLIDE = 10  

def create_four_slide_ppt(summary_data):
    prs = Presentation()

    # ===== Slide 1: Title =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "📊 FT.com Daily Summary"
    slide1.placeholders[1].text = f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}"
    slide1.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)

    # ===== Slide 2: Articles =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "📰 Articles"
    tf2 = slide2.placeholders[1].text_frame
    tf2.clear()
    for article in summary_data.get("articles", []):
        p = tf2.add_paragraph()
        p.text = f"- {article['title']}"
        p.font.size = Pt(12)

    # ===== Slide 3: Summaries =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "📝 Summaries"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    count = 0

    for article in summary_data.get("articles", []):
        if count >= MAX_PARAGRAPHS_PER_SLIDE:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "📝 Summaries (Continued)"
            tf = slide.placeholders[1].text_frame
            tf.clear()
            count = 0

        p_title = tf.add_paragraph()
        p_title.text = f"{article['title']}:"
        p_title.font.bold = True
        p_title.font.size = Pt(12)

        p_summary = tf.add_paragraph()
        p_summary.text = article.get("llm_summary", "")
        p_summary.font.size = Pt(10)

        count += 2

    # ===== Slide 4: References =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "🔗 References / Links"
    tf4 = slide4.placeholders[1].text_frame
    tf4.clear()
    for article in summary_data.get("articles", []):
        p = tf4.add_paragraph()
        p.text = f"- {article['link']}"
        p.font.size = Pt(10)

    # Save PPT
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_file = os.path.join(OUTPUT_DIR, f"FT_Summary_4Slides_{datetime.utcnow().date()}.pptx")
    prs.save(output_file)
    print(f" PPT saved: {output_file}")

if __name__ == "__main__":
    if not os.path.exists(INPUT_FILE):
        raise FileNotFoundError(f"Input file not found: {INPUT_FILE}")

    with open(INPUT_FILE, "r", encoding="utf-8") as f:
        summary_data = json.load(f)

    create_four_slide_ppt(summary_data)

